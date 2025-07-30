import os
import uuid
import subprocess
from mammoth import convert_to_html
from xhtml2pdf import pisa  # or use weasyprint if installed
from pdf2docx import Converter
from PIL import Image
from pptx import Presentation
import pandas as pd
import pdfplumber
import fitz
import zipfile
from xlsx2html import xlsx2html
from xhtml2pdf import pisa
import io
from pptx.util import Inches

def zip_files(file_paths, zip_path):
    with zipfile.ZipFile(zip_path, 'w') as zipf:
        for path in file_paths:
            arcname = os.path.basename(path)
            zipf.write(path, arcname)

# DOCX to PDF
def docx_to_pdf(input_path, output_path):
    try:
        subprocess.run(["unoconv", "-f", "pdf", "-o", output_path, input_path], check=True)
    except subprocess.CalledProcessError as e:
        raise RuntimeError(f"Failed to convert DOCX to PDF using unoconv: {e}")


# PDF to DOCX
def pdf_to_word(input_path, output_path):
    cv = Converter(input_path)
    cv.convert(output_path, start=0, end=None)
    cv.close()


# Images to PDF
def images_to_pdf(file_paths, output_path):
    images = [Image.open(f).convert("RGB") for f in file_paths]
    images[0].save(output_path, save_all=True, append_images=images[1:])


# PDF to Images
def pdf_to_images(file_bytes):
    os.makedirs("temp", exist_ok=True)
    pdf = fitz.open(stream=file_bytes, filetype="pdf")
    output_paths = []

    for i, page in enumerate(pdf):
        pix = page.get_pixmap(dpi=200)
        path = f"temp/page_{uuid.uuid4().hex}_{i+1}.png"
        pix.save(path)
        output_paths.append(path)

    pdf.close()
    return output_paths

# PPTX to PDF
def pptx_to_pdf(input_path, output_path):
    binary_path = os.path.join(os.path.dirname(__file__), "../bin/OfficeToPDF")
    binary_path = os.path.abspath(binary_path)

    try:
        subprocess.run([
            binary_path, input_path, output_path
        ], check=True)
    except subprocess.CalledProcessError as e:
        raise RuntimeError(f"PPTX to PDF conversion failed: {e}")

# PDF to PPTX
def pdf_to_pptx(input_path, output_path):
    doc = fitz.open(input_path)
    prs = Presentation()
    prs.slide_width = Inches(11.69)
    prs.slide_height = Inches(8.27)

    image_paths = []

    for i, page in enumerate(doc):
        pix = page.get_pixmap(dpi=150)
        img_path = f"temp_pdf_slide_{uuid.uuid4().hex}_{i+1}.png"
        pix.save(img_path)
        image_paths.append(img_path)

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

    prs.save(output_path)

    for path in image_paths:
        try:
            os.remove(path)
        except Exception as e:
            print(f"Warning: failed to delete temp image {path}: {e}")


# XLSX to PDF
def xlsx_to_pdf(input_path, output_path):
    html_stream = io.StringIO()
    xlsx2html(input_path, html_stream, sheet="Sheet1")
    html_content = html_stream.getvalue()

    with open(output_path, "wb") as pdf_file:
        pisa_status = pisa.CreatePDF(io.StringIO(html_content), dest=pdf_file)

    if pisa_status.err:
        raise Exception("Failed to convert Excel to PDF with formatting.")

# PDF to XLSX
def pdf_to_xlsx(input_path, output_path):
    all_data = []
    with pdfplumber.open(input_path) as pdf:
        for page in pdf.pages:
            table = page.extract_table()
            if table:
                df = pd.DataFrame(table[1:], columns=table[0])
                all_data.append(df)
    if all_data:
        result = pd.concat(all_data, ignore_index=True)
        result.to_excel(output_path, index=False)
    else:
        raise ValueError("No tables found in PDF.")
